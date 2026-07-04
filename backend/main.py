import sys
import os as _os

# Exe compilé sans console (PyInstaller console=False) : stdout/stderr sont None,
# ce qui fait planter uvicorn/logging (.isatty()). On les redirige vers le néant.
if sys.stdout is None:
    sys.stdout = open(_os.devnull, "w", encoding="utf-8")
if sys.stderr is None:
    sys.stderr = open(_os.devnull, "w", encoding="utf-8")

from fastapi import FastAPI, UploadFile, File, HTTPException, Depends, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from typing import List, Dict, Optional
import json
import time
import os
from dotenv import load_dotenv

load_dotenv()

from crag_engine import build_crag_graph, initial_state

app = FastAPI(title="KING2MO RAG API")

# CORS : uniquement les origines explicitement autorisées (M1).
# Configurable via ALLOWED_ORIGINS (valeurs séparées par des virgules).
_allowed = os.environ.get("ALLOWED_ORIGINS", "http://localhost:3000,http://127.0.0.1:3000,http://localhost:3050,http://127.0.0.1:3050")
ALLOWED_ORIGINS = [o.strip() for o in _allowed.split(",") if o.strip()]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Auth obligatoire : exige le BACKEND_API_TOKEN défini dans le .env
def require_token(x_api_token: Optional[str] = Header(default=None)):
    # Pas de mot de passe par défaut codé en dur : le .env doit le définir.
    expected = os.environ.get("BACKEND_API_TOKEN", "")
    if not expected:
        raise HTTPException(status_code=500, detail="BACKEND_API_TOKEN non configuré sur le serveur.")
    if x_api_token != expected:
        raise HTTPException(status_code=401, detail="Mot de passe d'accès invalide ou manquant.")
    return True

class ChatRequest(BaseModel):
    query: str = Field(..., min_length=1, max_length=10000)
    mode: str = "hybrid"
    history: List[Dict[str, str]] = [] # Ajout de l'historique
    image_base64: Optional[str] = None

@app.post("/api/chat")
async def chat(request: ChatRequest, _auth: bool = Depends(require_token)):
    gemini_key = os.environ.get("GEMINI_API_KEY", "")
    tavily_key = os.environ.get("TAVILY_API_KEY", "")
    graph = build_crag_graph(gemini_key, tavily_key)
    state = initial_state(request.query, request.history, request.mode, request.image_base64)

    # Générateur SYNCHRONE : FastAPI l'exécute dans un threadpool, donc la
    # génération LLM ne bloque plus la boucle d'événements du serveur.
    def event_stream():
        try:
            start_time = time.time()

            # Une seule exécution du graphe : on accumule l'état final
            # au fil des updates streamées (avant : stream + invoke = double exécution,
            # double latence et double coût LLM).
            final_state = dict(state)
            for update in graph.stream(state):
                for node, out in update.items():
                    if isinstance(out, dict):
                        final_state.update(out)
                    yield f"data: {json.dumps({'type': 'status', 'node': node})}\n\n"

            # Format documents (avec URL cliquable quand disponible)
            sources = []
            for doc in final_state.get("documents", []):
                meta = doc.metadata or {}
                source_name = meta.get("source", "local")
                url = meta.get("url") or (source_name if str(source_name).startswith("http") else None)
                snippet = doc.page_content[:300]
                if len(doc.page_content) > 300:
                    snippet += "..."
                sources.append({
                    "source": meta.get("title") or source_name,
                    "url": url,
                    "content": snippet,
                })

            result = {
                "type": "result",
                "generation": final_state.get("generation", ""),
                "sources": sources,
                "corrections": final_state.get("corrections", 0),
                "web_used": final_state.get("web_search") == "yes" or request.mode == "web",
                "duration": round(time.time() - start_time, 2),
                "input_tokens": final_state.get("input_tokens", 0),
                "output_tokens": final_state.get("output_tokens", 0),
                "search_count": final_state.get("search_count", 0),
            }
            yield f"data: {json.dumps(result)}\n\n"
        except Exception as e:
            error_msg = str(e)
            if "unknown variant image_url" in error_msg:
                error_msg = "❌ Le modèle d'IA actuel (ex: DeepSeek) ne supporte pas l'analyse d'images. Veuillez utiliser un modèle compatible (ex: Gemini 1.5, GPT-4o, Claude 3) pour envoyer des images."
            yield f"data: {json.dumps({'type': 'error', 'message': error_msg})}\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")

from fastapi import UploadFile, File
import tempfile
import uuid
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Limite de taille d'upload (C4) : 25 Mo par défaut, configurable.
MAX_UPLOAD_BYTES = int(os.environ.get("MAX_UPLOAD_MB", "25")) * 1024 * 1024

@app.post("/api/upload")
async def upload_document(file: UploadFile = File(...), _auth: bool = Depends(require_token)):
    # Validation du type côté serveur
    safe_name = os.path.basename(file.filename or "")  # neutralise le path traversal
    ext = os.path.splitext(safe_name)[1].lower()
    allowed_exts = [".pdf", ".txt", ".docx", ".xlsx", ".pptx"]
    
    if ext not in allowed_exts:
        return {"message": f"Seuls les fichiers {', '.join(allowed_exts)} sont acceptés.", "status": "error"}

    # Lecture bornée en taille (C4) pour éviter un déni de service.
    content = await file.read(MAX_UPLOAD_BYTES + 1)
    if len(content) > MAX_UPLOAD_BYTES:
        return {"message": f"Fichier trop volumineux (max {MAX_UPLOAD_BYTES // (1024*1024)} Mo).", "status": "error"}
    if not content:
        return {"message": "Fichier vide.", "status": "error"}

    # Nom temporaire unique et neutre (M6) — jamais dérivé du nom client.
    temp_file = os.path.join(tempfile.gettempdir(), f"upload_{uuid.uuid4().hex}{ext}")
    try:
        with open(temp_file, "wb") as buffer:
            buffer.write(content)

        docs = []
        if ext == ".pdf":
            loader = PyPDFLoader(temp_file)
            docs = loader.load()
        elif ext == ".txt":
            from langchain_core.documents import Document
            with open(temp_file, "r", encoding="utf-8") as f:
                docs = [Document(page_content=f.read())]
        elif ext == ".docx":
            from langchain_core.documents import Document
            import docx
            doc = docx.Document(temp_file)
            text = "\n".join([p.text for p in doc.paragraphs])
            docs = [Document(page_content=text)]
        elif ext == ".pptx":
            from langchain_core.documents import Document
            import pptx
            prs = pptx.Presentation(temp_file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            docs = [Document(page_content="\n".join(text_runs))]
        elif ext == ".xlsx":
            from langchain_core.documents import Document
            import openpyxl
            wb = openpyxl.load_workbook(temp_file, data_only=True)
            text_runs = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text_runs.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    row_texts = [str(cell) for cell in row if cell is not None]
                    if row_texts:
                        text_runs.append(" | ".join(row_texts))
            docs = [Document(page_content="\n".join(text_runs))]
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = text_splitter.split_documents(docs)

        # On force la métadonnée "source" au nom d'origine assaini (traçabilité + suppression).
        for d in splits:
            d.metadata["source"] = safe_name

        from crag_engine import _get_embeddings, CHROMA_DB_DIR
        from langchain_community.vectorstores import Chroma

        vectorstore = Chroma(
            persist_directory=CHROMA_DB_DIR,
            embedding_function=_get_embeddings(),
        )
        
        # Vérification si le document existe déjà (m3)
        db_data = vectorstore.get(include=["metadatas"])
        existing_sources = {meta.get("source") for meta in db_data.get("metadatas", [])}
        if safe_name in existing_sources:
            return {"message": f"Le document '{safe_name}' existe déjà dans la base.", "status": "warning"}
            
        vectorstore.add_documents(splits)

        return {"message": f"{len(splits)} segments ajoutés à la base.", "status": "success"}
    except Exception as e:
        return {"message": str(e), "status": "error"}
    finally:
        # Nettoyage systématique du fichier temporaire (M6).
        if os.path.exists(temp_file):
            try:
                os.remove(temp_file)
            except OSError:
                pass

@app.get("/api/documents")
async def get_documents(_auth: bool = Depends(require_token)):
    try:
        from crag_engine import _get_embeddings, CHROMA_DB_DIR
        from langchain_community.vectorstores import Chroma
        vectorstore = Chroma(persist_directory=CHROMA_DB_DIR, embedding_function=_get_embeddings())
        
        # On récupère tous les IDs et métadonnées pour extraire les noms de fichiers uniques
        db_data = vectorstore.get(include=["metadatas"])
        metadatas = db_data.get("metadatas", [])
        
        # Extraire les sources uniques (noms de fichiers)
        unique_sources = set()
        for meta in metadatas:
            if "source" in meta and not str(meta["source"]).startswith("http"):
                # Clean up temp prefix if present
                src = meta["source"]
                if src.startswith("temp_"):
                    src = src[5:]
                unique_sources.add(src)
                
        return {"documents": list(unique_sources), "status": "success"}
    except Exception as e:
        return {"documents": [], "status": "error", "message": str(e)}

@app.delete("/api/documents")
async def clear_all_documents(_auth: bool = Depends(require_token)):
    try:
        from crag_engine import _get_embeddings, CHROMA_DB_DIR
        from langchain_community.vectorstores import Chroma
        vectorstore = Chroma(persist_directory=CHROMA_DB_DIR, embedding_function=_get_embeddings())
        
        # Récupérer tous les IDs
        db_data = vectorstore.get()
        ids = db_data.get("ids", [])
        
        if ids:
            vectorstore.delete(ids=ids)
            
        return {"message": "Base de données vidée.", "status": "success"}
    except Exception as e:
        return {"message": str(e), "status": "error"}

@app.delete("/api/documents/{filename}")
async def delete_document(filename: str, _auth: bool = Depends(require_token)):
    try:
        from crag_engine import _get_embeddings, CHROMA_DB_DIR
        from langchain_community.vectorstores import Chroma
        vectorstore = Chroma(persist_directory=CHROMA_DB_DIR, embedding_function=_get_embeddings())
        
        db_data = vectorstore.get(include=["metadatas"])
        ids = db_data.get("ids", [])
        metadatas = db_data.get("metadatas", [])
        
        # Trouver les IDs liés à ce fichier (avec ou sans préfixe temp_)
        ids_to_delete = []
        for doc_id, meta in zip(ids, metadatas):
            src = meta.get("source", "")
            if src == filename or src == f"temp_{filename}":
                ids_to_delete.append(doc_id)
                
        if ids_to_delete:
            vectorstore.delete(ids=ids_to_delete)
            return {"message": f"Document '{filename}' supprimé ({len(ids_to_delete)} segments).", "status": "success"}
        else:
            from fastapi import HTTPException
            raise HTTPException(status_code=404, detail="Document introuvable.")
            
    except Exception as e:
        from fastapi import HTTPException
        raise HTTPException(status_code=500, detail=str(e))

# ---------------------------------------------------------------------
# Paramètres (fournisseur LLM + clés API), modifiables depuis l'interface.
# Écrits dans le .env à côté de l'exe (ou de main.py en dev).
# ---------------------------------------------------------------------
_VALID_PROVIDERS = {"", "gemini", "deepseek", "gemini-openai"}

def _env_file_path() -> str:
    import sys as _sys
    if getattr(_sys, "frozen", False):
        return os.path.join(os.path.dirname(_sys.executable), ".env")
    return os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")

def _mask_key(value: str) -> str:
    if not value:
        return ""
    return value[:4] + "…" + value[-4:] if len(value) > 8 else "•" * len(value)

class SettingsUpdate(BaseModel):
    llm_provider: Optional[str] = None
    gemini_api_key: Optional[str] = Field(default=None, max_length=300)
    tavily_api_key: Optional[str] = Field(default=None, max_length=300)

@app.get("/api/settings")
async def get_settings(_auth: bool = Depends(require_token)):
    return {
        "llm_provider": os.environ.get("LLM_PROVIDER", ""),
        "llm_api_key_masked": _mask_key(os.environ.get("GEMINI_API_KEY", "")),
        "tavily_api_key_masked": _mask_key(os.environ.get("TAVILY_API_KEY", "")),
    }

@app.post("/api/settings")
async def update_settings(update: SettingsUpdate, _auth: bool = Depends(require_token)):
    provider = (update.llm_provider or "").strip().lower()
    if provider not in _VALID_PROVIDERS:
        return {"status": "error", "message": f"Fournisseur inconnu : {provider}"}

    changes = {}
    if provider:
        changes["LLM_PROVIDER"] = provider
    if update.gemini_api_key and update.gemini_api_key.strip():
        changes["GEMINI_API_KEY"] = update.gemini_api_key.strip()
    if update.tavily_api_key and update.tavily_api_key.strip():
        changes["TAVILY_API_KEY"] = update.tavily_api_key.strip()
    if not changes:
        return {"status": "warning", "message": "Aucun changement fourni."}

    # Relit le .env existant en préservant l'ordre et les clés non modifiées
    env_path = _env_file_path()
    values, order = {}, []
    if os.path.exists(env_path):
        with open(env_path, "r", encoding="utf-8") as f:
            for raw in f.read().splitlines():
                s = raw.strip()
                if not s or s.startswith("#") or "=" not in s:
                    continue
                k, v = s.split("=", 1)
                k = k.strip()
                if k not in values:
                    order.append(k)
                values[k] = v.strip()

    for k, v in changes.items():
        values[k] = f'"{v}"'
        if k not in order:
            order.append(k)
        os.environ[k] = v  # effet immédiat, sans redémarrage

    with open(env_path, "w", encoding="utf-8") as f:
        for k in order:
            f.write(f"{k}={values[k]}\n")

    return {"status": "success", "message": "Paramètres enregistrés ✓"}

# ---------------------------------------------------------------------
# Serveur Static (Frontend Monolithe)
# ---------------------------------------------------------------------
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, HTMLResponse
import sys

# Résolution du chemin absolu du dossier "out" du frontend
base_dir = os.path.dirname(os.path.abspath(__file__))
frontend_build_path = os.path.join(base_dir, "..", "frontend", "out")

# Si on est dans un exécutable PyInstaller (_MEIPASS), le dossier est différent
if hasattr(sys, '_MEIPASS'):
    frontend_build_path = os.path.join(sys._MEIPASS, "frontend_out")

# Si le dossier de build existe (version compilée ou standalone), on le monte sur la racine
if os.path.isdir(frontend_build_path):
    app.mount("/_next", StaticFiles(directory=os.path.join(frontend_build_path, "_next")), name="next_assets")
    
    # On gère manuellement la racine et les 404 pour l'application mono-page Next.js
    @app.get("/")
    async def serve_index():
        return FileResponse(os.path.join(frontend_build_path, 'index.html'))

    @app.exception_handler(404)
    async def custom_404_handler(request, exc):
        if request.url.path.startswith("/api/"):
            return HTMLResponse(content="API Route Not Found", status_code=404)
        return FileResponse(os.path.join(frontend_build_path, 'index.html'))

if __name__ == "__main__":
    import uvicorn
    import threading
    import webbrowser
    import socket
    import traceback

    # Fichier de log à côté de l'exe (ou du script en mode dev)
    if getattr(sys, "frozen", False):
        _app_dir = os.path.dirname(sys.executable)
    else:
        _app_dir = os.path.dirname(os.path.abspath(__file__))
    LOG_FILE = os.path.join(_app_dir, "king2mo_error.log")

    def _port_in_use(port: int) -> bool:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            return s.connect_ex(("127.0.0.1", port)) == 0

    def _pick_free_port() -> int:
        """Essaie 8000 puis des alternatives ; sinon port attribué par l'OS.
        Sur Windows, des plages de ports peuvent être réservées (Hyper-V/WSL)
        et le bind sur 8000 échoue silencieusement."""
        for candidate in (8000, 8080, 8501, 3050):
            try:
                with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
                    s.bind(("127.0.0.1", candidate))
                return candidate
            except OSError:
                continue
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.bind(("127.0.0.1", 0))
            return s.getsockname()[1]

    PORT = _pick_free_port()
    URL = f"http://127.0.0.1:{PORT}"

    def _find_edge():
        import subprocess  # noqa: F401
        candidates = [
            os.path.expandvars(r"%ProgramFiles(x86)%\Microsoft\Edge\Application\msedge.exe"),
            os.path.expandvars(r"%ProgramFiles%\Microsoft\Edge\Application\msedge.exe"),
            os.path.expandvars(r"%LocalAppData%\Microsoft\Edge\Application\msedge.exe"),
        ]
        for p in candidates:
            if os.path.exists(p):
                return p
        return None

    def _open_app_window(block: bool):
        """Ouvre l'interface dans une fenêtre application (Edge --app).
        Si block=True, la fermeture de la fenêtre arrête l'application."""
        import subprocess
        import tempfile
        edge = _find_edge()
        if edge:
            # Profil persistant (pas dans %TEMP%) : conserve l'historique et
            # les réglages du navigateur intégré entre les sessions.
            profile = os.path.join(
                os.environ.get("LOCALAPPDATA", tempfile.gettempdir()),
                "KING2MO", "window_profile",
            )
            os.makedirs(profile, exist_ok=True)
            proc = subprocess.Popen([edge, f"--app={URL}", f"--user-data-dir={profile}"])
            if block:
                proc.wait()
                os._exit(0)  # fenêtre fermée -> on quitte le serveur
        else:
            webbrowser.open(URL)  # secours si Edge introuvable

    def _wait_server_then_open_window():
        for _ in range(240):  # jusqu'à 2 min (1er lancement plus lent)
            if _port_in_use(PORT):
                break
            time.sleep(0.5)
        _open_app_window(block=True)

    PORT_FILE = os.path.join(_app_dir, "king2mo.port")

    try:
        # Une instance tourne-t-elle déjà ? (port mémorisé au dernier lancement)
        try:
            with open(PORT_FILE, "r", encoding="utf-8") as f:
                previous_port = int(f.read().strip())
            if _port_in_use(previous_port):
                URL = f"http://127.0.0.1:{previous_port}"
                _open_app_window(block=False)
                sys.exit(0)
        except (OSError, ValueError):
            pass

        with open(PORT_FILE, "w", encoding="utf-8") as f:
            f.write(str(PORT))

        threading.Thread(target=_wait_server_then_open_window, daemon=True).start()

        # 127.0.0.1 : évite l'alerte du pare-feu Windows (0.0.0.0 la déclenche).
        # log_config=None : évite le formatter uvicorn incompatible sans console.
        uvicorn.run(app, host="127.0.0.1", port=PORT, log_config=None)
    except BaseException as e:
        # Capture aussi SystemExit(1) émis par uvicorn quand le démarrage échoue.
        if isinstance(e, SystemExit) and e.code in (0, None):
            raise
        with open(LOG_FILE, "a", encoding="utf-8") as f:
            f.write(traceback.format_exc() + "\n")
        raise
