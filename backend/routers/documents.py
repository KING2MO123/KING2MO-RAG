import os
import tempfile
import uuid
import re as _re
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form

from security import require_token

router = APIRouter()

# Limite de taille d'upload (C4) : 25 Mo par défaut, configurable.
MAX_UPLOAD_BYTES = int(os.environ.get("MAX_UPLOAD_MB", "25")) * 1024 * 1024
# (N3) Limite de taille décompressée pour les archives Office (anti zip-bomb).
MAX_DECOMPRESSED_BYTES = int(os.environ.get("MAX_DECOMPRESSED_MB", "200")) * 1024 * 1024

@router.post("/api/upload")
async def upload_document(file: UploadFile = File(...), scope: str = Form("global"), _auth: bool = Depends(require_token)):
    # (R-13) Portée du document : "global" (visible partout, défaut) ou un id
    # de conversation (32 hex). On valide le format pour éviter toute injection
    # dans les métadonnées.
    doc_scope = "global"
    if scope and scope != "global" and _re.fullmatch(r"[a-f0-9]{32}", scope):
        doc_scope = scope

    # Validation du type côté serveur
    safe_name = os.path.basename(file.filename or "")  # neutralise le path traversal
    ext = os.path.splitext(safe_name)[1].lower()
    allowed_exts = [".pdf", ".txt", ".docx", ".xlsx", ".pptx"]
    
    if ext not in allowed_exts:
        raise HTTPException(status_code=400, detail=f"Seuls les fichiers {', '.join(allowed_exts)} sont acceptés.")

    # Lecture bornée en taille (C4) pour éviter un déni de service.
    content = await file.read(MAX_UPLOAD_BYTES + 1)
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail=f"Fichier trop volumineux (max {MAX_UPLOAD_BYTES // (1024*1024)} Mo).")
    if not content:
        raise HTTPException(status_code=400, detail="Fichier vide.")

    # Nom temporaire unique et neutre (M6) — jamais dérivé du nom client.
    temp_file = os.path.join(tempfile.gettempdir(), f"upload_{uuid.uuid4().hex}{ext}")
    try:
        with open(temp_file, "wb") as buffer:
            buffer.write(content)

        # (N3) Anti zip-bomb : les fichiers Office sont des archives zip dont
        # la décompression peut saturer la mémoire. On vérifie la taille
        # décompressée annoncée avant tout parsing.
        if ext in (".docx", ".xlsx", ".pptx"):
            import zipfile
            try:
                with zipfile.ZipFile(temp_file) as zf:
                    total_uncompressed = sum(info.file_size for info in zf.infolist())
            except zipfile.BadZipFile:
                raise HTTPException(status_code=400, detail="Fichier Office invalide ou corrompu.")
            if total_uncompressed > MAX_DECOMPRESSED_BYTES:
                raise HTTPException(
                    status_code=413,
                    detail=f"Fichier Office refusé : taille décompressée excessive (max {MAX_DECOMPRESSED_BYTES // (1024*1024)} Mo).",
                )

        docs = []
        if ext == ".pdf":
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(temp_file)
            docs = loader.load()
        elif ext == ".txt":
            from langchain_core.documents import Document
            # (N7) Tolérance à l'encodage pour éviter les crashs sur les fichiers non-UTF-8
            with open(temp_file, "r", encoding="utf-8", errors="replace") as f:
                docs = [Document(page_content=f.read())]
        elif ext == ".docx":
            from langchain_core.documents import Document
            import docx
            doc = docx.Document(temp_file)
            text = "\n".join([p.text for p in doc.paragraphs])
            docs = [Document(page_content=text)]
        elif ext == ".pptx":
            from langchain_core.documents import Document
            import pptx
            prs = pptx.Presentation(temp_file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            docs = [Document(page_content="\n".join(text_runs))]
        elif ext == ".xlsx":
            from langchain_core.documents import Document
            import openpyxl
            wb = openpyxl.load_workbook(temp_file, data_only=True)
            text_runs = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text_runs.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    row_texts = [str(cell) for cell in row if cell is not None]
                    if row_texts:
                        text_runs.append(" | ".join(row_texts))
            docs = [Document(page_content="\n".join(text_runs))]
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        splits = text_splitter.split_documents(docs)

        # (R-5) Aucun texte exploitable : cas typique d'un PDF scanné (images
        # sans couche texte). Sans ce garde, l'utilisateur voyait « 0 segments »
        # sans comprendre pourquoi. On l'avertit explicitement.
        if not any((d.page_content or "").strip() for d in splits):
            hint = " (PDF probablement scanné : il ne contient que des images, sans texte reconnu)" if ext == ".pdf" else ""
            return {
                "message": f"Aucun texte n'a pu être extrait de « {safe_name} »{hint}. Rien n'a été ajouté.",
                "status": "warning",
            }

        # On force la métadonnée "source" au nom d'origine assaini (traçabilité + suppression).
        # (R-13) On tague aussi la portée (global ou id de conversation).
        for d in splits:
            d.metadata["source"] = safe_name
            d.metadata["scope"] = doc_scope

        from vectorstore import get_vectorstore
        vectorstore = get_vectorstore()
        
        # Vérification si le document existe déjà (m3)
        db_data = vectorstore.get(include=["metadatas"])
        existing_sources = {meta.get("source") for meta in db_data.get("metadatas", [])}
        if safe_name in existing_sources:
            return {"message": f"Le document '{safe_name}' existe déjà dans la base.", "status": "warning"}
            
        vectorstore.add_documents(splits)

        from vectorstore import invalidate_caches
        invalidate_caches()

        return {"message": f"{len(splits)} segments ajoutés à la base.", "status": "success"}
    except HTTPException:
        raise  # ne pas transformer les erreurs 4xx volontaires en 500
    except Exception as e:
        import logging
        logging.error(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail="Une erreur interne est survenue lors de l'upload.")
    finally:
        # Nettoyage systématique du fichier temporaire (M6).
        if os.path.exists(temp_file):
            try:
                os.remove(temp_file)
            except OSError:
                pass

@router.get("/api/documents")
async def get_documents(_auth: bool = Depends(require_token)):
    try:
        from vectorstore import get_vectorstore
        vectorstore = get_vectorstore()
        
        # On récupère tous les IDs et métadonnées pour extraire les noms de fichiers uniques
        db_data = vectorstore.get(include=["metadatas"])
        metadatas = db_data.get("metadatas", [])
        
        # Extraire les sources uniques (noms de fichiers)
        unique_sources = set()
        for meta in metadatas:
            if "source" in meta and not str(meta["source"]).startswith("http"):
                # Clean up temp prefix if present
                src = meta["source"]
                if src.startswith("temp_"):
                    src = src[5:]
                unique_sources.add(src)
                
        return {"documents": list(unique_sources), "status": "success"}
    except Exception as e:
        import logging
        logging.error(f"Get documents error: {e}")
        raise HTTPException(status_code=500, detail="Une erreur interne est survenue.")

@router.delete("/api/documents")
async def clear_all_documents(_auth: bool = Depends(require_token)):
    try:
        from vectorstore import get_vectorstore
        vectorstore = get_vectorstore()
        
        # Récupérer tous les IDs
        db_data = vectorstore.get()
        ids = db_data.get("ids", [])
        
        if ids:
            vectorstore.delete(ids=ids)
            from vectorstore import invalidate_caches
            invalidate_caches()

        return {"message": "Base de données vidée.", "status": "success"}
    except Exception as e:
        import logging
        logging.error(f"Clear documents error: {e}")
        raise HTTPException(status_code=500, detail="Une erreur interne est survenue.")

@router.delete("/api/documents/{filename}")
async def delete_document(filename: str, _auth: bool = Depends(require_token)):
    try:
        from vectorstore import get_vectorstore
        vectorstore = get_vectorstore()
        
        db_data = vectorstore.get(include=["metadatas"])
        ids = db_data.get("ids", [])
        metadatas = db_data.get("metadatas", [])
        
        # Trouver les IDs liés à ce fichier (avec ou sans préfixe temp_)
        ids_to_delete = []
        for doc_id, meta in zip(ids, metadatas):
            src = meta.get("source", "")
            if src == filename or src == f"temp_{filename}":
                ids_to_delete.append(doc_id)
                
        if ids_to_delete:
            vectorstore.delete(ids=ids_to_delete)
            from vectorstore import invalidate_caches
            invalidate_caches()
            return {"message": f"Document '{filename}' supprimé ({len(ids_to_delete)} segments).", "status": "success"}
        else:
            raise HTTPException(status_code=404, detail="Document introuvable.")
    except Exception as e:
        import logging
        logging.error(f"Delete document error: {str(e)}")
        if isinstance(e, HTTPException):
            raise
        raise HTTPException(status_code=500, detail="Une erreur interne est survenue.")
